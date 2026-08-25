from __future__ import annotations

import re
import uuid
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

from .config import settings


def safe_name(title: str) -> str:
    slug = re.sub(r"[^a-zA-Z0-9_-]+", "-", title.strip()).strip("-").lower() or "artifact"
    return f"{slug[:60]}-{uuid.uuid4().hex[:8]}"


def create_artifact(kind: str, title: str, content: str) -> Path:
    root = settings.data_dir / "artifacts"
    base = safe_name(title)
    lines = [line.strip() for line in content.splitlines() if line.strip()]
    if kind == "docx":
        path = root / f"{base}.docx"
        doc = Document()
        doc.add_heading(title, 0)
        for line in lines:
            if line.startswith("#"):
                doc.add_heading(line.lstrip("# "), level=min(line.count("#"), 3))
            else:
                doc.add_paragraph(line)
        doc.save(path)
    elif kind == "pptx":
        path = root / f"{base}.pptx"
        prs = Presentation()
        chunks = [lines[i : i + 5] for i in range(0, len(lines), 5)] or [[content]]
        for idx, chunk in enumerate(chunks[:12]):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title if idx == 0 else chunk[0][:80]
            body = slide.placeholders[1].text_frame
            body.text = "\n".join(chunk[1:] if idx else chunk)[:1800]
        prs.save(path)
    elif kind == "xlsx":
        path = root / f"{base}.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = "Nexus AI"
        ws.append([title])
        for i, line in enumerate(lines, start=2):
            ws.cell(i, 1, line)
        ws.column_dimensions["A"].width = 100
        wb.save(path)
    elif kind == "pdf":
        path = root / f"{base}.pdf"
        pdf = canvas.Canvas(str(path), pagesize=A4)
        width, height = A4
        y = height - 60
        pdf.setFont("Helvetica-Bold", 18)
        pdf.drawString(50, y, title[:75])
        y -= 35
        pdf.setFont("Helvetica", 10)
        for line in lines:
            for start in range(0, len(line), 95):
                if y < 60:
                    pdf.showPage()
                    pdf.setFont("Helvetica", 10)
                    y = height - 60
                pdf.drawString(50, y, line[start : start + 95])
                y -= 14
        pdf.save()
    else:
        path = root / f"{base}.md"
        path.write_text(f"# {title}\n\n{content}", encoding="utf-8")
    return path
